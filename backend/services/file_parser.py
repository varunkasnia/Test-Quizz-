import io
import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract
from docx import Document
from typing import Optional


def parse_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF file."""
    try:
        pdf_file = io.BytesIO(file_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        
        if not text.strip():
            raise ValueError("No text could be extracted from the PDF")
        
        return text.strip()
    
    except Exception as e:
        raise Exception(f"Failed to parse PDF: {str(e)}")


def parse_text(file_bytes: bytes) -> str:
    """Extract text from TXT file."""
    try:
        # Try UTF-8 first, then fallback to latin-1
        try:
            text = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            text = file_bytes.decode('latin-1')
        
        if not text.strip():
            raise ValueError("The text file is empty")
        
        return text.strip()
    
    except Exception as e:
        raise Exception(f"Failed to parse text file: {str(e)}")


def parse_pptx(file_bytes: bytes) -> str:
    """Extract text from PowerPoint file."""
    try:
        pptx_file = io.BytesIO(file_bytes)
        presentation = Presentation(pptx_file)
        
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        if not text.strip():
            raise ValueError("No text could be extracted from the PowerPoint file")
        
        return text.strip()
    
    except Exception as e:
        raise Exception(f"Failed to parse PowerPoint: {str(e)}")


def parse_docx(file_bytes: bytes) -> str:
    """Extract text from Word document."""
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = Document(docx_file)
        
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        if not text.strip():
            raise ValueError("No text could be extracted from the Word document")
        
        return text.strip()
    
    except Exception as e:
        raise Exception(f"Failed to parse Word document: {str(e)}")


def parse_image(file_bytes: bytes) -> str:
    """Extract text from image using OCR."""
    try:
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        
        if not text.strip():
            raise ValueError("No text could be extracted from the image. Make sure the image contains readable text.")
        
        return text.strip()
    
    except Exception as e:
        raise Exception(f"Failed to parse image: {str(e)}")


def parse_file(filename: str, file_bytes: bytes) -> str:
    """
    Parse file based on extension and return extracted text.
    Supports: PDF, TXT, PPTX, DOCX, and images (PNG, JPG, JPEG)
    """
    extension = filename.lower().split('.')[-1]
    
    parsers = {
        'pdf': parse_pdf,
        'txt': parse_text,
        'pptx': parse_pptx,
        'ppt': parse_pptx,
        'docx': parse_docx,
        'doc': parse_docx,
        'png': parse_image,
        'jpg': parse_image,
        'jpeg': parse_image,
    }
    
    parser = parsers.get(extension)
    if not parser:
        raise ValueError(
            f"Unsupported file type: {extension}. "
            f"Supported types: {', '.join(parsers.keys())}"
        )
    
    return parser(file_bytes)
